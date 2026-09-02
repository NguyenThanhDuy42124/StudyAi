"""
Parse PPTX → plain text dùng python-pptx.
Xử lý: slide titles, text boxes, speaker notes.
"""
from loguru import logger
from pptx import Presentation
from pptx.util import Pt


def parse_pptx(file_path: str) -> str:
    """
    Extract toàn bộ text từ file PPTX.
    Bao gồm: titles, text boxes, notes của từng slide.
    
    Args:
        file_path: Đường dẫn đến file PPTX
    
    Returns:
        Text đã extract, có đánh số slide
    """
    try:
        prs = Presentation(file_path)
        slides_text: list[str] = []
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts: list[str] = [f"[Slide {slide_num}]"]
            
            # Extract text từ shapes (title, content, text boxes)
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
            
            # Extract speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"[Ghi chú: {notes_text}]")
            
            if len(slide_parts) > 1:  # Có nội dung ngoài header
                slides_text.append("\n".join(slide_parts))
        
        result = "\n\n".join(slides_text)
        logger.info(f"PPTX parse OK: {len(prs.slides)} slides, {len(result)} ký tự")
        return result
        
    except Exception as e:
        logger.error(f"Lỗi parse PPTX {file_path}: {e}")
        raise
